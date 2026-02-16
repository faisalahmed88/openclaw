#!/usr/bin/env python3
"""Generate AI Basics & OpenClaw Architecture PowerPoint presentation."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Theme colors ──────────────────────────────────────────────────────
BG_DARK   = RGBColor(0x1E, 0x1E, 0x2E)   # dark bg
BG_CARD   = RGBColor(0x2A, 0x2A, 0x3C)   # card bg
ACCENT    = RGBColor(0x7C, 0x3A, 0xED)   # purple accent
ACCENT2   = RGBColor(0x06, 0xB6, 0xD4)   # cyan accent
TEXT_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
TEXT_LIGHT = RGBColor(0xCC, 0xCC, 0xDD)
TEXT_MUTED = RGBColor(0x99, 0x99, 0xAA)
ORANGE     = RGBColor(0xF9, 0x73, 0x16)
GREEN      = RGBColor(0x22, 0xC5, 0x5E)
RED        = RGBColor(0xEF, 0x44, 0x44)
YELLOW     = RGBColor(0xFA, 0xCC, 0x15)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

# ── Helpers ───────────────────────────────────────────────────────────
def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape_rect(slide, left, top, w, h, fill_color, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    return shape

def add_text_box(slide, left, top, w, h, text, font_size=18, color=TEXT_WHITE, bold=False, alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = alignment
    return tf

def add_bullet_slide(slide, left, top, w, h, items, font_size=16, color=TEXT_LIGHT, spacing=Pt(6)):
    txBox = slide.shapes.add_textbox(left, top, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.space_after = spacing
        p.level = 0
    return tf

def add_card(slide, left, top, w, h, title, bullets, title_color=ACCENT2, bullet_color=TEXT_LIGHT, title_size=20, bullet_size=14):
    add_shape_rect(slide, left, top, w, h, BG_CARD, border_color=RGBColor(0x44, 0x44, 0x55))
    add_text_box(slide, left + Inches(0.2), top + Inches(0.15), w - Inches(0.4), Inches(0.5), title, font_size=title_size, color=title_color, bold=True)
    add_bullet_slide(slide, left + Inches(0.2), top + Inches(0.6), w - Inches(0.4), h - Inches(0.8), bullets, font_size=bullet_size, color=bullet_color)

def title_slide(title_text, subtitle_text=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide, BG_DARK)
    add_text_box(slide, Inches(1), Inches(2.2), Inches(11), Inches(1.5), title_text, font_size=44, color=TEXT_WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    if subtitle_text:
        add_text_box(slide, Inches(1), Inches(3.8), Inches(11), Inches(1), subtitle_text, font_size=22, color=TEXT_MUTED, alignment=PP_ALIGN.CENTER)
    return slide

def section_slide(title_text, subtitle_text=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    # accent bar
    add_shape_rect(slide, Inches(1), Inches(3.05), Inches(3), Inches(0.08), ACCENT)
    add_text_box(slide, Inches(1), Inches(1.8), Inches(11), Inches(1.2), title_text, font_size=40, color=TEXT_WHITE, bold=True)
    if subtitle_text:
        add_text_box(slide, Inches(1), Inches(3.4), Inches(11), Inches(1), subtitle_text, font_size=20, color=TEXT_MUTED)
    return slide

def content_slide(title_text):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    add_shape_rect(slide, Inches(0), Inches(0), prs.slide_width, Inches(1.1), RGBColor(0x25, 0x25, 0x38))
    add_text_box(slide, Inches(0.8), Inches(0.2), Inches(11), Inches(0.7), title_text, font_size=30, color=TEXT_WHITE, bold=True)
    return slide


# =====================================================================
#  SLIDE 1 — Title
# =====================================================================
s = title_slide(
    "AI Fundamentals & OpenClaw Architecture",
    "Models · Agents · Tokens · MCP · Sessions · Skills — and how OpenClaw ties it all together"
)
add_text_box(s, Inches(1), Inches(5.2), Inches(11), Inches(0.6),
    "Faisal Ahmed  ·  February 2026  ·  Programming Context (VS Code / Anthropic / OpenAI Codex)",
    font_size=16, color=TEXT_MUTED, alignment=PP_ALIGN.CENTER)

# =====================================================================
#  SLIDE 2 — Agenda
# =====================================================================
s = content_slide("Agenda")
items = [
    "1.  AI Basics — what is AI, ML, LLMs",
    "2.  Models — types, open-source vs paid, capabilities comparison",
    "3.  Anthropic Claude vs OpenAI GPT/Codex — key differences",
    "4.  Tokens & Context Windows — how models consume input",
    "5.  Agents — what they are, multi-agent orchestration, agent handoff",
    "6.  Configurable Tools & MCP (Model Context Protocol)",
    "7.  Sessions, Context Management, Add-Context",
    "8.  Custom Agents & Skills (deep dive)",
    "9.  VS Code as the AI IDE — Copilot, extensions, ACP",
    "10. OpenClaw — architecture walkthrough using all the above concepts",
]
add_bullet_slide(s, Inches(0.8), Inches(1.4), Inches(11), Inches(5.5), items, font_size=20, color=TEXT_LIGHT, spacing=Pt(10))

# =====================================================================
#  SLIDE 3 — AI Basics
# =====================================================================
s = section_slide("1. AI Basics", "Artificial Intelligence → Machine Learning → Deep Learning → LLMs")

# =====================================================================
#  SLIDE 4 — AI Basics detail
# =====================================================================
s = content_slide("AI Basics — Core Concepts")
add_card(s, Inches(0.5), Inches(1.4), Inches(3.8), Inches(5.4),
    "Artificial Intelligence",
    ["Broad field: machines performing tasks that typically require human intelligence",
     "Includes: vision, NLP, robotics, reasoning, planning",
     "Sub-fields: ML, deep learning, reinforcement learning"])
add_card(s, Inches(4.7), Inches(1.4), Inches(3.8), Inches(5.4),
    "Machine Learning",
    ["Algorithms that learn patterns from data",
     "Supervised: labeled data (classification, regression)",
     "Unsupervised: clustering, dimensionality reduction",
     "Reinforcement: trial-and-error with reward signals"])
add_card(s, Inches(8.9), Inches(1.4), Inches(3.8), Inches(5.4),
    "Large Language Models (LLMs)",
    ["Neural networks trained on massive text corpora",
     "Predict next token given context (autoregressive)",
     "Capabilities: text generation, code, reasoning, translation",
     "Examples: GPT-4o, Claude Opus, Llama 3, Gemini"])

# =====================================================================
#  SLIDE 5 — Models Overview
# =====================================================================
s = section_slide("2. Models", "Types · Open-Source vs Paid · Capabilities")

# =====================================================================
#  SLIDE 6 — Model Types
# =====================================================================
s = content_slide("Model Types")
add_card(s, Inches(0.5), Inches(1.4), Inches(5.8), Inches(2.4),
    "By Architecture",
    ["Autoregressive (GPT, Claude, Llama) — generate token-by-token left→right",
     "Encoder-only (BERT) — bidirectional, good for classification/embeddings",
     "Encoder-Decoder (T5, BART) — seq2seq tasks (translation, summarization)",
     "Mixture of Experts (MoE) — Mixtral, GPT-4 — route tokens to specialist sub-networks"])
add_card(s, Inches(6.7), Inches(1.4), Inches(5.8), Inches(2.4),
    "By Modality",
    ["Text-only: GPT-4o-mini, Claude Haiku",
     "Multimodal: GPT-4o (text+image+audio), Claude Opus (text+image)",
     "Code-specialized: Codex, DeepSeek Coder, CodeLlama",
     "Embedding: text-embedding-3, voyage-3 — vector representations"])
add_card(s, Inches(0.5), Inches(4.2), Inches(5.8), Inches(2.8),
    "Open-Source (Free)",
    ["Llama 3.1 (Meta) — 8B/70B/405B params, very capable",
     "Mistral / Mixtral (Mistral AI) — efficient MoE",
     "DeepSeek V3/R1 — strong reasoning, code",
     "Qwen 2.5 (Alibaba) — multilingual",
     "Phi-4 (Microsoft) — small but strong",
     "Run locally via Ollama, llama.cpp, vLLM"])
add_card(s, Inches(6.7), Inches(4.2), Inches(5.8), Inches(2.8),
    "Paid / API",
    ["OpenAI: GPT-4o, GPT-4o-mini, o1/o3 (reasoning), Codex",
     "Anthropic: Claude Opus 4.6, Sonnet 4.5, Haiku 3.5",
     "Google: Gemini 2.5 Pro/Flash",
     "AWS Bedrock: hosts Claude, Llama, Cohere, Titan",
     "Pricing: per input/output token (varies 10x–100x by model)",
     "OAuth subscription: Claude Pro/Max, ChatGPT Plus/Pro"])

# =====================================================================
#  SLIDE 7 — Anthropic vs OpenAI
# =====================================================================
s = content_slide("Anthropic Claude vs OpenAI GPT / Codex")
add_card(s, Inches(0.5), Inches(1.4), Inches(5.8), Inches(5.4),
    "Anthropic Claude",
    ["Models: Opus 4.6 (flagship), Sonnet 4.5, Haiku 3.5",
     "Strengths: long-context (200K tokens), safety-focused, strong reasoning",
     "Extended thinking: dedicated thinking budget for complex tasks",
     "Better prompt-injection resistance (Constitutional AI)",
     "Code: excellent at TypeScript, Python, multi-file refactoring",
     "API: Messages API with system/user/assistant roles",
     "Auth: OAuth (Pro/Max subscription) or API key",
     "VS Code: powers GitHub Copilot (Claude model option)"],
    title_color=ORANGE)
add_card(s, Inches(6.7), Inches(1.4), Inches(5.8), Inches(5.4),
    "OpenAI GPT / Codex",
    ["Models: GPT-4o (multimodal), o1/o3 (reasoning), Codex (code agent)",
     "Strengths: multimodal (text+image+audio), function calling, JSON mode",
     "Codex: cloud sandbox agent for autonomous code tasks",
     "Broader ecosystem: ChatGPT, DALL-E, Whisper, TTS",
     "API: Chat Completions + Assistants API + Responses API",
     "Context: 128K tokens (GPT-4o), 200K (o1)",
     "Auth: API key or ChatGPT Plus/Pro subscription",
     "VS Code: GitHub Copilot (GPT-4o option), Codex agent"],
    title_color=GREEN)

# =====================================================================
#  SLIDE 8 — Tokens & Context Window
# =====================================================================
s = section_slide("3. Tokens & Context Window", "The currency of LLMs")

# =====================================================================
#  SLIDE 9 — Tokens detail
# =====================================================================
s = content_slide("Tokens — How Models Read Text")
add_card(s, Inches(0.5), Inches(1.4), Inches(5.8), Inches(5.4),
    "What Are Tokens?",
    ["Atomic units after tokenization (subwords/pieces)",
     "\"Hello world\" → [\"Hello\", \" world\"] = 2 tokens",
     "\"Anthropic\" → [\"Anthrop\", \"ic\"] = 2 tokens",
     "Rule of thumb: 1 token ≈ 4 characters ≈ 0.75 words",
     "1K tokens ≈ 750 English words",
     "Code tokens are different: symbols, brackets count as tokens",
     "Tokenizers: BPE (OpenAI), SentencePiece (LLama/Claude)",
     "Each model has its own tokenizer — token counts differ"])
add_card(s, Inches(6.7), Inches(1.4), Inches(5.8), Inches(5.4),
    "Why Tokens Matter",
    ["Cost: billed per input + output tokens",
     "  GPT-4o: ~$2.50/1M input, ~$10/1M output",
     "  Claude Opus: ~$15/1M input, ~$75/1M output",
     "  Claude Haiku: ~$0.25/1M input, ~$1.25/1M output",
     "Speed: more tokens = slower generation",
     "Context window = max tokens model can see at once",
     "Attention complexity: roughly O(n²) with token count",
     "Prompt engineering: keep instructions concise to save tokens"])

# =====================================================================
#  SLIDE 10 — Context Window
# =====================================================================
s = content_slide("Context Window — Model Memory")
add_card(s, Inches(0.5), Inches(1.4), Inches(11.8), Inches(2.2),
    "What Is the Context Window?",
    ["Max tokens (prompt + completion) the model can attend to in one request",
     "Anything beyond the limit is truncated (usually from the start) — model forgets it",
     "Larger window = more context, but higher cost + latency"])
add_card(s, Inches(0.5), Inches(3.9), Inches(5.8), Inches(3.2),
    "Context Window Sizes",
    ["GPT-4o: 128K tokens (~96K words)",
     "GPT-4o-mini: 128K tokens",
     "o1 / o3: 200K tokens",
     "Claude Opus 4.6: 200K tokens",
     "Claude Sonnet 4.5: 200K tokens",
     "Gemini 2.5 Pro: 1M tokens",
     "Llama 3.1 405B: 128K tokens"])
add_card(s, Inches(6.7), Inches(3.9), Inches(5.8), Inches(3.2),
    "Best Practices",
    ["Put critical instructions at the START of the prompt",
     "Use retrieval (RAG): fetch relevant chunks, not full docs",
     "Summarize long histories instead of including verbatim",
     "Use embeddings + vector DB for long-term memory",
     "Chunk large files and process incrementally",
     "Monitor token usage to control costs"])

# =====================================================================
#  SLIDE 11 — Agents Section
# =====================================================================
s = section_slide("4. Agents", "Autonomous workers that use models + tools to accomplish tasks")

# =====================================================================
#  SLIDE 12 — What Is an Agent
# =====================================================================
s = content_slide("What Is an Agent?")
add_card(s, Inches(0.5), Inches(1.4), Inches(5.8), Inches(5.4),
    "Agent = Model + Tools + Loop",
    ["An agent is software that reasons, plans, and acts autonomously",
     "Uses an LLM as its 'brain' for reasoning",
     "Has access to tools (shell, HTTP, file I/O, APIs)",
     "Runs in a loop: Think → Act → Observe → Repeat",
     "Maintains state across steps (session, memory)",
     "Can ask clarifying questions or delegate to sub-agents",
     "Enforces policies: permissions, timeouts, safety guards"],
    title_color=ACCENT)
add_card(s, Inches(6.7), Inches(1.4), Inches(5.8), Inches(5.4),
    "Agent vs Model",
    ["Model: stateless prediction engine (input→output)",
     "Agent: stateful orchestrator (plan→execute→iterate)",
     "",
     "Model: single inference call",
     "Agent: multiple calls in a reasoning loop",
     "",
     "Model: no tools (just text/tokens)",
     "Agent: can call tools, read files, run code, browse",
     "",
     "Model: you control the loop",
     "Agent: controls its own loop (autonomous)"],
    title_color=ACCENT2)

# =====================================================================
#  SLIDE 13 — Multi-Agent
# =====================================================================
s = content_slide("Multi-Agent Orchestration")
add_card(s, Inches(0.5), Inches(1.4), Inches(3.8), Inches(5.4),
    "Parallel Agents",
    ["Multiple agents run concurrently on separate tasks",
     "Example: Agent A refactors code, Agent B writes tests, Agent C updates docs",
     "Each has its own session, tools, and working directory",
     "No shared state — communicate via files, messages, or queues",
     "OpenClaw: multiple CLI agents in separate terminals",
     "VS Code: Copilot multi-agent with workspace isolation"])
add_card(s, Inches(4.7), Inches(1.4), Inches(3.8), Inches(5.4),
    "Agent Communication",
    ["Agents can talk to each other via:",
     "• Shared files (git repo as shared state)",
     "• Message queues / event bus",
     "• Gateway WebSocket events (OpenClaw)",
     "• Agent Client Protocol (ACP) sessions",
     "• Tool calls that invoke other agents",
     "Pattern: coordinator agent delegates sub-tasks to worker agents"])
add_card(s, Inches(8.9), Inches(1.4), Inches(3.8), Inches(5.4),
    "Agent Handoff",
    ["Transfer a conversation/task from one agent to another",
     "Use cases:",
     "• Escalation: simple agent → expert agent",
     "• Specialization: code agent → review agent",
     "• Capacity: busy agent → available agent",
     "Handoff includes: session state, context, pending tasks",
     "OpenClaw: session key routing between agents"])

# =====================================================================
#  SLIDE 14 — Tools & MCP
# =====================================================================
s = section_slide("5. Configurable Tools & MCP", "Model Context Protocol — the USB-C of AI tools")

# =====================================================================
#  SLIDE 15 — Tools detail
# =====================================================================
s = content_slide("Configurable Tools")
add_card(s, Inches(0.5), Inches(1.4), Inches(5.8), Inches(5.4),
    "What Are Tools?",
    ["Functions/APIs an agent can call during reasoning",
     "Defined as JSON schemas (name, description, parameters)",
     "Agent decides WHEN and HOW to call them",
     "",
     "Common tool types:",
     "• Shell / terminal execution (bash, cmd)",
     "• File system (read, write, search, list)",
     "• HTTP / API calls (fetch URLs, REST APIs)",
     "• Browser automation (Playwright)",
     "• Database queries (SQL, vector search)",
     "• Custom plugins (any function you define)"])
add_card(s, Inches(6.7), Inches(1.4), Inches(5.8), Inches(5.4),
    "Tool Configuration",
    ["Tools are registered with the agent at startup",
     "Each tool has: name, description, input schema, handler",
     "Permission levels: auto-approve (read), prompt (write/exec)",
     "Tool categories:",
     "  read / search → safe, auto-approved",
     "  edit / write → requires confirmation",
     "  execute / bash → dangerous, always prompt",
     "  delete / move → destructive, always prompt",
     "Tools can be added/removed per session",
     "OpenClaw: tools configured via gateway + plugin SDK"])

# =====================================================================
#  SLIDE 16 — MCP
# =====================================================================
s = content_slide("MCP — Model Context Protocol")
add_card(s, Inches(0.5), Inches(1.4), Inches(11.8), Inches(2.0),
    "What Is MCP?",
    ["Open standard (by Anthropic) for connecting AI agents to external tools & data sources",
     "Think of it as 'USB-C for AI' — one protocol, many tools",
     "JSON-RPC over stdio or HTTP — transport-agnostic, language-agnostic"],
    title_color=ACCENT)
add_card(s, Inches(0.5), Inches(3.7), Inches(3.8), Inches(3.4),
    "MCP Architecture",
    ["MCP Host: the AI app (VS Code, Claude Desktop)",
     "MCP Client: manages connection to servers",
     "MCP Server: exposes tools + resources",
     "Transport: stdio (local) or HTTP/SSE (remote)",
     "Protocol: initialize → list tools → call tool"])
add_card(s, Inches(4.7), Inches(3.7), Inches(3.8), Inches(3.4),
    "MCP Capabilities",
    ["Tools: functions the agent can invoke",
     "Resources: data sources the agent can read",
     "Prompts: reusable prompt templates",
     "Sampling: let server request LLM completions",
     "Roots: file system roots for sandboxing",
     "Logging: structured diagnostic output"])
add_card(s, Inches(8.9), Inches(3.7), Inches(3.8), Inches(3.4),
    "MCP in Practice",
    ["VS Code: Copilot supports MCP servers",
     "Claude Desktop: native MCP integration",
     "OpenClaw: gateway acts as MCP-like tool hub",
     "Examples: GitHub MCP, Postgres MCP, Filesystem MCP",
     "Community: 1000s of MCP servers on GitHub",
     "Config: JSON in settings / mcp.json"])

# =====================================================================
#  SLIDE 17 — Sessions Section
# =====================================================================
s = section_slide("6. Sessions & Context Management", "How agents maintain state across interactions")

# =====================================================================
#  SLIDE 18 — Sessions detail
# =====================================================================
s = content_slide("Sessions, Context, Add-Context")
add_card(s, Inches(0.5), Inches(1.4), Inches(3.8), Inches(5.4),
    "Sessions",
    ["A session = persistent conversation state",
     "Contains: message history, tool results, metadata",
     "Identified by session key (e.g. 'agent:main:main')",
     "Can be: created, loaded, listed, reset, compacted",
     "Session store: in-memory, file-based, or DB-backed",
     "OpenClaw: sessions managed via gateway WebSocket",
     "Multiple sessions can run in parallel"])
add_card(s, Inches(4.7), Inches(1.4), Inches(3.8), Inches(5.4),
    "Context Management",
    ["Context = everything the model sees in a request",
     "Includes: system prompt + history + tools + user input",
     "Must fit within the context window",
     "Strategies when context is full:",
     "• Truncate oldest messages",
     "• Summarize/compact history",
     "• Use retrieval (RAG) for relevant chunks",
     "• Sliding window over conversation",
     "OpenClaw: auto-compaction via /compact command"])
add_card(s, Inches(8.9), Inches(1.4), Inches(3.8), Inches(5.4),
    "Add-Context",
    ["Dynamically inject context into agent prompts",
     "Types of added context:",
     "• Files: attach source code or docs",
     "• URLs: fetch and include web pages",
     "• Images: screenshots, diagrams",
     "• Resources: MCP resource URIs",
     "• Clipboard: paste content directly",
     "VS Code Copilot: #file, #selection, #codebase",
     "OpenClaw: attachments in chat.send, CWD prefix"])

# =====================================================================
#  SLIDE 19 — Custom Agents & Skills Section
# =====================================================================
s = section_slide("7. Custom Agents & Skills", "Building specialized agents with reusable capabilities")

# =====================================================================
#  SLIDE 20 — Custom Agents
# =====================================================================
s = content_slide("Custom Agents")
add_card(s, Inches(0.5), Inches(1.4), Inches(5.8), Inches(5.4),
    "What Are Custom Agents?",
    ["Agents configured for specific roles or domains",
     "Customized via: system prompt, tool set, model choice, policies",
     "",
     "Examples:",
     "• Code Review Agent — reads PRs, uses lint/test tools",
     "• DevOps Agent — manages deployments, uses SSH/kubectl",
     "• Docs Agent — writes documentation, uses Mintlify tools",
     "• Security Agent — audits code, uses SAST/DAST scanners",
     "",
     "Configuration approaches:",
     "• System prompt templates (persona + rules + examples)",
     "• Tool allow-lists (only give relevant tools)",
     "• Model selection (fast model for simple, powerful for complex)"])
add_card(s, Inches(6.7), Inches(1.4), Inches(5.8), Inches(5.4),
    "Agent Configuration in OpenClaw",
    ["System prompt: configurable per session/channel",
     "Model: selectable via /model command or config",
     "Thinking level: off/minimal/low/medium/high/xhigh",
     "Tools: registered via plugin SDK + gateway",
     "Permissions: auto-approve safe, prompt for dangerous",
     "Session key: routes to specific agent instance",
     "Channels: agent responds on WhatsApp, Discord, Telegram, etc.",
     "Hooks: lifecycle callbacks (before/after reply, on error)",
     "Cron: scheduled agent tasks"])

# =====================================================================
#  SLIDE 21 — Skills Deep Dive
# =====================================================================
s = content_slide("Skills — Deep Dive")
add_card(s, Inches(0.5), Inches(1.4), Inches(11.8), Inches(1.6),
    "What Are Skills?",
    ["Skills = reusable, structured knowledge documents that teach an agent HOW to do a specific task",
     "They are NOT code — they are markdown instructions that get injected into the agent's context when relevant",
     "Think of skills as 'playbooks' or 'runbooks' the agent follows step-by-step"],
    title_color=ACCENT)
add_card(s, Inches(0.5), Inches(3.3), Inches(3.8), Inches(3.8),
    "Skill Structure",
    ["Each skill contains:",
     "• Name: identifier (e.g. 'review-pr')",
     "• Description: when to use this skill",
     "• File: markdown with detailed instructions",
     "",
     "Skill file includes:",
     "• Step-by-step workflow",
     "• Tool usage patterns",
     "• Decision trees / conditionals",
     "• Output format templates",
     "• Error handling guidance"])
add_card(s, Inches(4.7), Inches(3.3), Inches(3.8), Inches(3.8),
    "How Skills Are Used",
    ["1. User asks agent to do something",
     "2. Agent (or system) matches task to a skill",
     "3. Skill instructions are injected into context",
     "4. Agent follows the skill step-by-step",
     "5. Skill guides: which tools to call, in what order",
     "",
     "Activation: keyword match, explicit command, or auto-detect",
     "Skills can reference other skills (chaining)",
     "Skills are version-controlled (in repo)"])
add_card(s, Inches(8.9), Inches(3.3), Inches(3.8), Inches(3.8),
    "OpenClaw Skills (examples)",
    ["review-pr: deterministic PR review workflow",
     "  → reads diff, runs checks, produces findings",
     "",
     "prepare-pr: push-safe PR preparation",
     "  → resolves findings, formats commits, pushes",
     "",
     "merge-pr: squash-merge with check gating",
     "  → verifies CI, pins SHA, merges",
     "",
     "mintlify: docs site maintenance",
     "  → follows Mintlify conventions, checks links",
     "",
     "Located in: .agents/skills/<name>/SKILL.md"])

# =====================================================================
#  SLIDE 22 — VS Code Section
# =====================================================================
s = section_slide("8. VS Code — The AI IDE", "GitHub Copilot · Extensions · Agent Client Protocol")

# =====================================================================
#  SLIDE 23 — VS Code detail
# =====================================================================
s = content_slide("VS Code as the AI Development Platform")
add_card(s, Inches(0.5), Inches(1.4), Inches(3.8), Inches(5.4),
    "GitHub Copilot",
    ["AI pair programmer built into VS Code",
     "Inline completions (ghost text)",
     "Chat panel: ask questions about code",
     "Agent mode: autonomous multi-step coding",
     "Models: GPT-4o, Claude Opus/Sonnet, Gemini",
     "MCP support: connect external tool servers",
     "Custom instructions: .github/copilot-instructions.md"])
add_card(s, Inches(4.7), Inches(1.4), Inches(3.8), Inches(5.4),
    "Context Variables (#)",
    ["#file — include a specific file",
     "#selection — current editor selection",
     "#codebase — semantic search across project",
     "#terminalLastCommand — last terminal output",
     "#problems — VS Code diagnostics",
     "#fetch — fetch a URL",
     "#changes — git diff",
     "These inject context into the prompt dynamically"])
add_card(s, Inches(8.9), Inches(1.4), Inches(3.8), Inches(5.4),
    "ACP (Agent Client Protocol)",
    ["Standard for IDE ↔ Agent communication",
     "Used by VS Code, JetBrains, terminal clients",
     "JSON-RPC over stdio (ndjson streams)",
     "Operations: initialize, newSession, prompt, cancel",
     "Session updates: streaming text, tool calls, commands",
     "Permission model: auto-approve safe, prompt dangerous",
     "OpenClaw: full ACP server + client in src/acp/"])

# =====================================================================
#  SLIDE 24 — Model Comparison
# =====================================================================
s = content_slide("Model Capabilities Comparison")
# Table-like layout using cards
add_card(s, Inches(0.5), Inches(1.4), Inches(3.0), Inches(1.4),
    "Claude Opus 4.6", ["Best: long-context reasoning, code refactoring, safety", "Context: 200K tokens"], title_color=ORANGE, bullet_size=13)
add_card(s, Inches(3.8), Inches(1.4), Inches(3.0), Inches(1.4),
    "Claude Sonnet 4.5", ["Best: balanced speed + quality, daily coding", "Context: 200K tokens"], title_color=ORANGE, bullet_size=13)
add_card(s, Inches(7.1), Inches(1.4), Inches(3.0), Inches(1.4),
    "Claude Haiku 3.5", ["Best: fast, cheap, simple tasks, classification", "Context: 200K tokens"], title_color=ORANGE, bullet_size=13)
add_card(s, Inches(10.4), Inches(1.4), Inches(2.5), Inches(1.4),
    "GPT-4o", ["Best: multimodal, function calling, JSON", "Context: 128K tokens"], title_color=GREEN, bullet_size=13)
add_card(s, Inches(0.5), Inches(3.1), Inches(3.0), Inches(1.4),
    "GPT o1 / o3", ["Best: complex reasoning, math, science", "Context: 200K tokens"], title_color=GREEN, bullet_size=13)
add_card(s, Inches(3.8), Inches(3.1), Inches(3.0), Inches(1.4),
    "Codex (OpenAI)", ["Best: autonomous code tasks in sandbox", "Context: 192K tokens"], title_color=GREEN, bullet_size=13)
add_card(s, Inches(7.1), Inches(3.1), Inches(3.0), Inches(1.4),
    "Gemini 2.5 Pro", ["Best: massive context, multimodal", "Context: 1M tokens"], title_color=ACCENT2, bullet_size=13)
add_card(s, Inches(10.4), Inches(3.1), Inches(2.5), Inches(1.4),
    "Llama 3.1 405B", ["Best: open-source, self-hosted, privacy", "Context: 128K tokens"], title_color=RED, bullet_size=13)
add_card(s, Inches(0.5), Inches(4.8), Inches(3.0), Inches(1.4),
    "DeepSeek R1", ["Best: open reasoning, code, low cost", "Context: 128K tokens"], title_color=RED, bullet_size=13)
add_card(s, Inches(3.8), Inches(4.8), Inches(3.0), Inches(1.4),
    "Mistral / Mixtral", ["Best: efficient MoE, EU-hosted, multilingual", "Context: 128K tokens"], title_color=RED, bullet_size=13)
add_card(s, Inches(7.1), Inches(4.8), Inches(5.8), Inches(1.4),
    "Key Takeaway",
    ["No single 'best' model — choose based on: task complexity, speed, cost, context needs, privacy",
     "OpenClaw supports ALL of these via provider/model configuration"],
    title_color=YELLOW, bullet_size=14)

# =====================================================================
#  SLIDE 25 — OpenClaw Section
# =====================================================================
s = section_slide("9. OpenClaw — Putting It All Together",
    "Multi-channel AI gateway with extensible messaging integrations")

# =====================================================================
#  SLIDE 26 — OpenClaw What Is It
# =====================================================================
s = content_slide("What Is OpenClaw?")
add_card(s, Inches(0.5), Inches(1.4), Inches(11.8), Inches(1.6),
    "Personal AI Assistant You Run on Your Own Devices",
    ["Open-source (MIT) · TypeScript/ESM · Node 22+ · CLI + Gateway + Multi-platform apps",
     "Answers you on WhatsApp, Telegram, Slack, Discord, Signal, iMessage, Teams, WebChat + 15 more channels",
     "The gateway is the control plane — the product is the assistant"],
    title_color=ACCENT)
add_card(s, Inches(0.5), Inches(3.3), Inches(5.8), Inches(3.8),
    "Core Architecture",
    ["Gateway Server: WebSocket + HTTP, manages sessions, routing, auth",
     "Agents: AI reasoning loop (models + tools + session state)",
     "Channels: message adapters (Telegram, Discord, Slack, etc.)",
     "Providers: LLM integrations (Anthropic, OpenAI, Bedrock, local)",
     "Plugin SDK: extension API for community channel/tool plugins",
     "CLI: Commander-based, interactive setup wizard (openclaw onboard)",
     "ACP Server: IDE integration via Agent Client Protocol",
     "Apps: macOS (Swift), iOS (Swift), Android (Kotlin)"])
add_card(s, Inches(6.7), Inches(3.3), Inches(5.8), Inches(3.8),
    "Key Stats (from audit)",
    ["1,875 TypeScript source files · 325,346 LOC",
     "50 directories under src/",
     "36 extensions (channel + tool plugins)",
     "1,226 test files (vitest + V8 coverage)",
     "326 English docs (Mintlify)",
     "9 CI workflows (Linux/Windows/macOS, Node+Bun)",
     "46 runtime dependencies",
     "Daily releases with dense fix lists"])

# =====================================================================
#  SLIDE 27 — OpenClaw Architecture Mapped to Concepts
# =====================================================================
s = content_slide("OpenClaw — Concepts in Action")
add_card(s, Inches(0.5), Inches(1.4), Inches(3.8), Inches(2.4),
    "🧠  Models (src/providers/)",
    ["Anthropic, OpenAI, Bedrock, local LLMs",
     "Model selection: /model command",
     "Failover: OAuth → API key rotation",
     "Thinking: off → xhigh (extended)"],
    title_color=ACCENT2, bullet_size=13)
add_card(s, Inches(4.7), Inches(1.4), Inches(3.8), Inches(2.4),
    "🤖  Agents (src/agents/)",
    ["Agent loop: prompt → reason → tool → respond",
     "Embedded runner, sandbox, sub-agents",
     "ACP gateway agent (src/acp/)",
     "Multi-agent: parallel CLI sessions"],
    title_color=ACCENT2, bullet_size=13)
add_card(s, Inches(8.9), Inches(1.4), Inches(3.8), Inches(2.4),
    "🔧  Tools (src/agents/, plugins/)",
    ["Shell exec, file I/O, browser (Playwright)",
     "HTTP/API calls, media processing",
     "Permission model: safe auto-approve",
     "Plugin-registered tools via SDK"],
    title_color=ACCENT2, bullet_size=13)
add_card(s, Inches(0.5), Inches(4.1), Inches(3.8), Inches(2.4),
    "📡  Channels (src/channels/, ext/)",
    ["Built-in: Telegram, Discord, Slack, Signal, etc.",
     "Extensions: Matrix, Nostr, IRC, Twitch, Teams",
     "Plugin SDK: build your own channel",
     "Allowlists, onboarding, typing indicators"],
    title_color=ACCENT2, bullet_size=13)
add_card(s, Inches(4.7), Inches(4.1), Inches(3.8), Inches(2.4),
    "💬  Sessions (src/sessions/, acp/)",
    ["Session keys: agent:main:main",
     "Create, load, list, reset, compact",
     "In-memory + gateway-managed stores",
     "Context management: auto-compaction"],
    title_color=ACCENT2, bullet_size=13)
add_card(s, Inches(8.9), Inches(4.1), Inches(3.8), Inches(2.4),
    "📚  Skills (.agents/skills/)",
    ["review-pr, prepare-pr, merge-pr, mintlify",
     "Markdown playbooks injected into context",
     "Step-by-step agent workflows",
     "Version-controlled, composable"],
    title_color=ACCENT2, bullet_size=13)

# =====================================================================
#  SLIDE 28 — OpenClaw src/ Directory Map
# =====================================================================
s = content_slide("OpenClaw — Source Directory Map")
add_card(s, Inches(0.3), Inches(1.4), Inches(4.1), Inches(5.5),
    "Core Infrastructure",
    ["src/gateway/ — WS server, HTTP API, protocol",
     "src/config/ — config loading, Zod schemas",
     "src/infra/ — heartbeat, state migrations",
     "src/routing/ — message routing, session keys",
     "src/cli/ — CLI wiring (Commander)",
     "src/commands/ — CLI commands (auth, config…)",
     "src/daemon/ — background process mgmt",
     "src/logging/ — log transports",
     "src/security/ — audit, verification"],
    bullet_size=12)
add_card(s, Inches(4.6), Inches(1.4), Inches(4.1), Inches(5.5),
    "AI & Agent Layer",
    ["src/agents/ — agent loop, tools, sandbox",
     "src/providers/ — LLM integrations",
     "src/acp/ — Agent Client Protocol",
     "src/auto-reply/ — reply chunking, streaming",
     "src/memory/ — QMD manager, sync",
     "src/sessions/ — session management",
     "src/browser/ — Playwright browser tool",
     "src/tts/ — text-to-speech",
     "src/media/ — MIME detection, storage"],
    bullet_size=12)
add_card(s, Inches(8.9), Inches(1.4), Inches(4.1), Inches(5.5),
    "Channels & Extensions",
    ["src/telegram/ — Telegram (grammy)",
     "src/discord/ — Discord (@buape/carbon)",
     "src/slack/ — Slack (@slack/bolt)",
     "src/signal/ — Signal",
     "src/imessage/ — iMessage",
     "src/whatsapp/ — WhatsApp (Baileys)",
     "src/channels/ — shared plugin system",
     "src/plugin-sdk/ — extension API",
     "extensions/ — 36 community plugins"],
    bullet_size=12)

# =====================================================================
#  SLIDE 29 — OpenClaw Data Flow
# =====================================================================
s = content_slide("OpenClaw — Message Flow (End to End)")
add_card(s, Inches(0.5), Inches(1.4), Inches(11.8), Inches(5.2),
    "User Message → AI Response (step by step)",
    ["1. USER sends message on a channel (e.g. WhatsApp, Discord, Telegram, CLI)",
     "",
     "2. CHANNEL ADAPTER receives the message, normalizes it, checks allowlist",
     "   → src/telegram/bot-handlers.ts, src/discord/monitor/, src/slack/, etc.",
     "",
     "3. ROUTING resolves the session key and routes to the correct agent session",
     "   → src/routing/ — session key resolution, channel → session mapping",
     "",
     "4. GATEWAY dispatches the message to the agent via WebSocket event frames",
     "   → src/gateway/server/ — ws-connection, message-handler",
     "",
     "5. AGENT receives the prompt, injects session context + system prompt + tools",
     "   → src/agents/ — pi-embedded-runner, context assembly, tool registration",
     "",
     "6. MODEL processes the prompt (Anthropic/OpenAI/etc.) and returns tokens",
     "   → src/providers/ — API calls, streaming, token counting",
     "",
     "7. AGENT executes any tool calls (shell, file, HTTP, plugin tools), feeds results back to model",
     "   → src/agents/bash-tools.exec.ts, src/browser/, plugin SDK tool handlers",
     "",
     "8. RESPONSE streams back through gateway → channel adapter → user's chat",
     "   → src/auto-reply/ — chunking, coalesce, queue → channel send functions"],
    title_color=ACCENT, bullet_size=14)

# =====================================================================
#  SLIDE 30 — Improvements
# =====================================================================
s = content_slide("OpenClaw — Fork Improvement Opportunities")
add_card(s, Inches(0.5), Inches(1.4), Inches(5.8), Inches(5.4),
    "High Impact — Quick Wins",
    ["14/36 extensions have ZERO tests (discord, signal, slack…)",
     "  → Add basic unit tests for message parsing & config",
     "",
     "23/36 extensions lack a README",
     "  → Add minimal docs (purpose, config, setup)",
     "",
     "Most extensions missing from docs navigation",
     "  → Update docs/docs.json Mintlify nav",
     "",
     "Plugin SDK has no getting-started guide",
     "  → Create docs/plugins/authoring.md"],
    title_color=GREEN)
add_card(s, Inches(6.7), Inches(1.4), Inches(5.8), Inches(5.4),
    "High Impact — Medium Effort",
    ["130 files exceed 500-LOC guideline (top: 1,156 LOC)",
     "  → Refactor: extract helpers, split modules",
     "",
     "Coverage exclusions hide real test gaps",
     "  → Gradually remove exclusions, add tests",
     "",
     "Dual schema (Zod + TypeBox) adds cognitive load",
     "  → Standardize on one (likely Zod v4)",
     "",
     "iOS CI is disabled",
     "  → Re-enable build verification",
     "",
     "Pre-release deps (node-pty, sqlite-vec, rolldown)",
     "  → Track stable releases and upgrade"],
    title_color=ORANGE)

# =====================================================================
#  SLIDE 31 — Summary
# =====================================================================
s = content_slide("Summary — Key Takeaways")
add_card(s, Inches(0.5), Inches(1.4), Inches(11.8), Inches(5.2),
    "Everything Connects",
    ["MODELS are the prediction engines — choose by task, cost, context window, capabilities",
     "",
     "TOKENS are the currency — every word costs tokens; context window limits what the model sees",
     "",
     "AGENTS are the orchestrators — they use models + tools in a reasoning loop to accomplish tasks",
     "",
     "TOOLS extend agents — file I/O, shell, HTTP, browser, plugins; configured and permission-gated",
     "",
     "MCP standardizes tool integration — 'USB-C for AI'; one protocol connects any tool to any agent",
     "",
     "SESSIONS maintain state — conversation history, context management, compaction, handoff",
     "",
     "SKILLS are agent playbooks — markdown instructions injected into context for step-by-step workflows",
     "",
     "OPENCLAW ties it all together — a self-hosted AI gateway that routes messages across 20+ channels,",
     "  orchestrates agents with configurable models/tools/skills, and integrates with IDEs via ACP"],
    title_color=ACCENT, bullet_size=16)

# =====================================================================
#  SLIDE 32 — Thank You
# =====================================================================
s = title_slide("Thank You", "Questions?")
add_text_box(s, Inches(1), Inches(5.0), Inches(11), Inches(0.5),
    "GitHub: github.com/faisalahmed88/openclaw  ·  Docs: docs.openclaw.ai",
    font_size=16, color=TEXT_MUTED, alignment=PP_ALIGN.CENTER)

# ── Save ──────────────────────────────────────────────────────────────
output_path = "/home/event/openclaw/AI_Basics_and_OpenClaw_Architecture.pptx"
prs.save(output_path)
print(f"✅ Presentation saved to {output_path}")
print(f"   Slides: {len(prs.slides)}")
